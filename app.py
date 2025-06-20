import os
import io
import json
import random
import difflib
from datetime import datetime
from pathlib import Path
from typing import Dict, Any

import streamlit as st
import openai
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import matplotlib.pyplot as plt
from PIL import Image

# Local .env for development; Streamlit Cloud uses st.secrets
from dotenv import load_dotenv
load_dotenv(Path(__file__).parent / ".env")

# ── Page config ─────────────────────────────────────────────────
st.set_page_config(layout="wide", page_title="Impact Project Room")

# ── OpenAI API key ────────────────────────────────────────────────
try:
    openai.api_key = st.secrets["OPENAI_API_KEY"]
except KeyError:
    st.error(
        "🔒 OPENAI_API_KEY missing. Add it under 'Manage app → Settings → Secrets'."
    )
    st.stop()

# ── Logo ──────────────────────────────────────────────────────────
logo_path = Path(__file__).with_name("logo.png")
if logo_path.exists():
    col1, col2 = st.columns([9, 1])
    with col2:
        st.image(Image.open(logo_path), width=60)

# ── Constants ─────────────────────────────────────────────────────
UPLOAD_FOLDER = "submissions"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

AI_FIELDS = [
    "Project Name", "Specific Sector(s)", "Region of operation", "Main country of current operations",
    "Business Model", "Maturity stage", "Core team", "Key risks",
    "Last 12 months revenues (USD)", "Breakeven year", "Market size or SOM (USD)",
    "Expected IRR (%)", "Financing need or round size (USD)", "Instrument",
    "Use of proceeds (%)", "Impact Area", "3 main SDGs targeted",
    "Problem", "Solution", "Barrier(s) to entry",
]
REVIEW_STAGES = [
    "Identified", "Intro call", "NDA and Deck", "Financials", "4-pager",
    "IC1", "IC2", "Local DD", "Raised", "Operating", "Exited", "Bankrupt",
]
SDG_OPTIONS = [
    "No poverty (SDG 1)", "Zero hunger (SDG 2)", "Good health and well-being (SDG 3)",
    "Quality education (SDG 4)", "Gender equality (SDG 5)", "Clean water and sanitation (SDG 6)",
    "Affordable and clean energy (SDG 7)", "Decent work and economic growth (SDG 8)",
    "Industry, innovation and infrastructure (SDG 9)", "Reduced inequalities (SDG 10)",
    "Sustainable cities and communities (SDG 11)", "Responsible consumption and production (SDG 12)",
    "Climate action (SDG 13)", "Life below water (SDG 14)", "Life on land (SDG 15)",
    "Peace, justice, and strong institutions (SDG 16)", "Partnerships for the goals (SDG 17)",
]
MATURITY_STAGES = ["Ideation", "Validation", "Pilot", "Growth", "Scale", "Mature"]
SECTOR_OPTIONS = [
    "Agriculture", "Air", "Biodiversity & ecosystems", "Climate", "Diversity & inclusion",
    "Education", "Employment / Livelihoods creation", "Energy", "Financial services",
    "Health", "Infrastructure", "Land", "Oceans & coastal zones",
    "Sustainable cities", "Sustainable consumption & production", "Sustainable tourism",
    "Water Treatment", "Other",
]
# ISO‑style country display list 
COUNTRY_OPTIONS = [
    "Afghanistan","Albania","Algeria","American Samoa","Andorra","Angola",
    "Anguilla","Antarctica","Antigua And Barbuda","Argentina","Armenia",
    "Aruba","Australia","Austria","Azerbaijan","Bahamas The","Bahrain",
    "Bangladesh","Barbados","Belarus","Belgium","Belize","Benin","Bermuda",
    "Bhutan","Bolivia","Bosnia and Herzegovina","Botswana","Bouvet Island",
    "Brazil","British Indian Ocean Territory","Brunei","Bulgaria",
    "Burkina Faso","Burundi","Cambodia","Cameroon","Canada","Cape Verde",
    "Cayman Islands","Central African Republic","Chad","Chile","China",
    "Christmas Island","Cocos (Keeling) Islands","Colombia","Comoros",
    "Republic Of The Congo","Democratic Republic Of The Congo","Cook Islands",
    "Costa Rica","Cote D'Ivoire (Ivory Coast)","Croatia (Hrvatska)","Cuba",
    "Cyprus","Czech Republic","Denmark","Djibouti","Dominica",
    "Dominican Republic","East Timor","Ecuador","Egypt","El Salvador",
    "Equatorial Guinea","Eritrea","Estonia","Ethiopia",
    "External Territories of Australia","Falkland Islands","Faroe Islands",
    "Fiji Islands","Finland","France","French Guiana","French Polynesia",
    "French Southern Territories","Gabon","Gambia The","Georgia","Germany",
    "Ghana","Gibraltar","Greece","Greenland","Grenada","Guadeloupe","Guam",
    "Guatemala","Guernsey and Alderney","Guinea","Guinea-Bissau","Guyana",
    "Haiti","Heard and McDonald Islands","Honduras","Hong Kong S.A.R.",
    "Hungary","Iceland","India","Indonesia","Iran","Iraq","Ireland","Israel",
    "Italy","Jamaica","Japan","Jersey","Jordan","Kazakhstan","Kenya",
    "Kiribati","Korea North","Korea South","Kuwait","Kyrgyzstan","Laos",
    "Latvia","Lebanon","Lesotho","Liberia","Libya","Liechtenstein",
    "Lithuania","Luxembourg","Macau S.A.R.","Macedonia","Madagascar",
    "Malawi","Malaysia","Maldives","Mali","Malta","Marshall Islands",
    "Martinique","Mauritania","Mauritius","Mayotte","Mexico","Micronesia",
    "Moldova","Monaco","Mongolia","Montenegro","Montserrat","Morocco",
    "Mozambique","Myanmar (Burma)","Namibia","Nauru","Nepal","Netherlands",
    "Netherlands Antilles","New Caledonia","New Zealand","Nicaragua","Niger",
    "Nigeria","Niue","Norfolk Island","Northern Mariana Islands","Norway",
    "Oman","Pakistan","Palau","Palestinian Territories","Panama",
    "Papua New Guinea","Paraguay","Peru","Philippines","Pitcairn Islands",
    "Poland","Portugal","Puerto Rico","Qatar","Reunion","Romania","Russia",
    "Rwanda","Saint Helena","Saint Kitts and Nevis","Saint Lucia",
    "Saint Pierre and Miquelon","Saint Vincent and the Grenadines","Samoa",
    "San Marino","Sao Tome and Principe","Saudi Arabia","Senegal","Serbia",
    "Seychelles","Sierra Leone","Singapore","Slovakia","Slovenia",
    "Solomon Islands","Somalia","South Africa",
    "South Georgia and the South Sandwich Islands","Spain","Sri Lanka",
    "Sudan","Suriname","Svalbard and Jan Mayen","Swaziland","Sweden",
    "Switzerland","Syria","Taiwan","Tajikistan","Tanzania","Thailand",
    "Timor-Leste (East Timor)","Togo","Tokelau","Tonga","Trinidad and Tobago",
    "Tunisia","Turkey","Turkmenistan","Turks and Caicos Islands","Tuvalu",
    "Uganda","Ukraine","United Arab Emirates","United Kingdom",
    "United States","United States Minor Outlying Islands","Uruguay",
    "Uzbekistan","Vanuatu","Vatican City","Venezuela","Vietnam",
    "Virgin Islands, British","Virgin Islands, U.S.","Wallis and Futuna",
    "Western Sahara","Yemen","Zambia","Zimbabwe",
]

# ── Standard options for review-stage input ───────────────────────
SDG_OPTIONS = [
    "No poverty (SDG 1)",
    "Zero hunger (SDG 2)",
    "Good health and well-being (SDG 3)",
    "Quality education (SDG 4)",
    "Gender equality (SDG 5)",
    "Clean water and sanitation (SDG 6)",
    "Affordable and clean energy (SDG 7)",
    "Decent work and economic growth (SDG 8)",
    "Industry, innovation and infrastructure (SDG 9)",
    "Reduced inequalities (SDG 10)",
    "Sustainable cities and communities (SDG 11)",
    "Responsible consumption and production (SDG 12)",
    "Climate action (SDG 13)",
    "Life below water (SDG 14)",
    "Life on land (SDG 15)",
    "Peace, justice, and strong institutions (SDG 16)",
    "Partnerships for the goals (SDG 17)",
]

MATURITY_STAGES = [
    "Ideation",
    "Validation",
    "Pilot",
    "Growth",
    "Scale",
    "Mature",
]

# ── Helpers ────────────────────────────────────────────────────────
def _match_sdgs(raw_list):
    matched = []
    for s in raw_list:
        cands = difflib.get_close_matches(s, SDG_OPTIONS, n=1, cutoff=0.6)
        if cands and cands[0] not in matched:
            matched.append(cands[0])
        if len(matched) >= 3:
            break
    return matched

def extract_text_from_pdf(path: str) -> str:
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf": return extract_text_from_pdf(path)
    if ext == ".docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if ext == ".pptx":
        prs = Presentation(path)
        return "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    if ext in (".xls", ".xlsx"):
        try: return pd.read_excel(path).to_csv(index=False)
        except: return ""
    return ""

def _parse_json_from_string(payload: str) -> Dict[str, Any]:
    try: return json.loads(payload)
    except:
        if '{' in payload and '}' in payload:
            snippet = payload[payload.find('{'):payload.rfind('}')+1]
            try: return json.loads(snippet)
            except: pass
    return {}

def summarize_project_with_gpt(full_text: str) -> Dict[str, Any]:
    text = full_text[:15000]
    system_prompt = (
        "You are an expert impact investment analyst. "
        "Extract these fields and return only a JSON object with keys: "
        + ", ".join(AI_FIELDS)
    )
    user_prompt = f"Pitch Content:\n{text}"
    try:
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",
            messages=[
                {"role":"system","content":system_prompt},
                {"role":"user","content":user_prompt},
            ],
            temperature=0.0,
            max_tokens=2000,
        )
        raw = resp.choices[0].message.content.strip()
        summary = _parse_json_from_string(raw)
    except Exception as e:
        st.error(f"OpenAI API error: {e}")
        summary = {}
    for k in AI_FIELDS:
        summary.setdefault(k, "Unknown")
    return summary

def render_summary_grid(summary: Dict[str, Any]):
    cols = st.columns(3)
    for i, field in enumerate(AI_FIELDS):
        val = summary.get(field, "–")
        if isinstance(val, list):
            val = "; ".join(str(x) for x in val)
        cols[i%3].markdown(f"**{field}**  \
{val}")

# ── Helper: save the submission to disk ───────────────────────────
def _save_submission(meta: dict, files, summary: dict):
    """
    Creates a timestamped folder in /submissions and writes:
    • each uploaded file
    • info.txt with human‑readable metadata
    • summary_gpt.txt containing the AI JSON
    """
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_proj = meta["Project"].replace(" ", "_")
    fld = os.path.join(UPLOAD_FOLDER, f"{safe_proj}_{ts}")
    os.makedirs(fld, exist_ok=True)

    # Save uploaded files
    for upl in files:
        fpath = os.path.join(fld, upl.name)
        with open(fpath, "wb") as f:
            f.write(upl.read())

    # Write metadata
    meta_text = "\n".join(f"{k}: {v}" for k, v in meta.items()) + "\nNDA: Accepted\n"
    with open(os.path.join(fld, "info.txt"), "w", encoding="utf-8") as mf:
        mf.write(meta_text)

    # Save AI summary
    with open(os.path.join(fld, "summary_gpt.txt"), "w", encoding="utf-8") as sf:
        json.dump(summary, sf, ensure_ascii=False, indent=2)

    # Mirror submission folder to SharePoint
    try:
        _upload_to_sharepoint(fld, os.path.basename(fld))
    except Exception as e:
        print("SharePoint upload error:", e)

    # Generate and save edit PIN for the entrepreneur
    pin = f"{random.randint(0, 9999):04d}"
    cred_path = os.path.join(fld, "credentials.json")
    with open(cred_path, "w", encoding="utf-8") as cf:
        json.dump({"pin": pin}, cf)
    return fld, pin


# ── Helper: update an existing submission ────────────────────────
def _update_submission(folder: str, meta: dict, summary: dict):
    """Overwrite metadata and summary in an existing submission folder."""
    base = os.path.join(UPLOAD_FOLDER, folder)
    # Overwrite metadata
    meta_text = "\n".join(f"{k}: {v}" for k, v in meta.items()) + "\nNDA: Accepted\n"
    with open(os.path.join(base, "info.txt"), "w", encoding="utf-8") as mf:
        mf.write(meta_text)
    # Overwrite summary
    with open(os.path.join(base, "summary_gpt.txt"), "w", encoding="utf-8") as sf:
        json.dump(summary, sf, ensure_ascii=False, indent=2)

def _upload_to_sharepoint(local_folder: str, project_folder: str):
    """
    Upload all files from local_folder to a SharePoint document library,
    creating a subfolder named project_folder.
    """
    # Pull SharePoint config from secrets
    site_url      = st.secrets["SP_SITE_URL"]
    client_id     = st.secrets["SP_CLIENT_ID"]
    client_secret = st.secrets["SP_CLIENT_SECRET"]
    library       = st.secrets["SP_DOC_LIBRARY"]  # e.g. "Shared Documents/ImpactSubmissions"

    # Authenticate
    creds = ClientCredential(client_id, client_secret)
    ctx   = ClientContext(site_url).with_credentials(creds)

    # Create (or get) the project folder
    root_folder      = ctx.web.get_folder_by_server_relative_url(library)
    project_sp_folder = root_folder.add_folder(project_folder).execute_query()

    # Upload each file in the local folder
    for fname in os.listdir(local_folder):
        path = os.path.join(local_folder, fname)
        with open(path, "rb") as f:
            content = f.read()
        project_sp_folder.upload_file(fname, content).execute_query()

# ── Helper: rerun compatible with all Streamlit versions ──────────
def _rerun():
    if hasattr(st, "rerun"):            # Streamlit ≥ 1.27
        st.rerun()
    else:                               # older versions
        st.experimental_rerun()

# ── Helper: send e‑mail alert when a new submission arrives ──────
import smtplib, ssl
from email.message import EmailMessage

def email_admin(subject: str, body: str):
    """Send a plain‑text e‑mail to the address set in secrets or fallback literal."""
    host = st.secrets.get("SMTP_HOST")
    port = int(st.secrets.get("SMTP_PORT", 0))
    user = st.secrets.get("SMTP_USER")
    pwd  = st.secrets.get("SMTP_PASS")
    to   = st.secrets.get("ADMIN_EMAIL", "NCM@kickimpact.com")

    if not all([host, port, user, pwd]):
        print("Email credentials missing: alert not sent.")
        return

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"]    = user
    msg["To"]      = to
    msg.set_content(body)

    ctx = ssl.create_default_context()
    try:
        if port == 465:   # SSL
            with smtplib.SMTP_SSL(host, port, context=ctx) as srv:
                srv.login(user, pwd)
                srv.send_message(msg)
        else:             # STARTTLS
            with smtplib.SMTP(host, port) as srv:
                srv.starttls(context=ctx)
                srv.login(user, pwd)
                srv.send_message(msg)
    except Exception as e:
        print("E‑mail error:", e)

# ── Routing ───────────────────────────────────────────────────────
is_admin = str(st.query_params.get("adminNCM", "")).lower() == "true"

# ═════════════════════════════════════════════════════════════════
# ADMIN DASHBOARD
# ═════════════════════════════════════════════════════════════════
if is_admin:
    st.title("🛠️ Admin Dashboard")
    # Sidebar filter categories
    st.sidebar.header("Filters")
    st.sidebar.subheader("Categorical Variables")
    # Admin filters
    hq_filter = st.sidebar.text_input("Headquarters country contains")
    sector_filter = st.sidebar.multiselect("Sector", SECTOR_OPTIONS)
    main_country_filter = st.sidebar.text_input("Main country of current operations contains")
    sdg_filter = st.sidebar.multiselect("3 main SDGs targeted", SDG_OPTIONS)
    geography_filter = st.sidebar.multiselect("Region of operation", [
        "Global", "Western Economies", "Africa", "Asia", "SEA", "Latam"])
    maturity_filter = st.sidebar.multiselect("Maturity stage", MATURITY_STAGES)

    st.sidebar.markdown("---")
    st.sidebar.subheader("Numerical Variables")

    # Numeric filters
    min_rev = st.sidebar.number_input("Min last 12 months revenues (USD)", min_value=0, value=0, step=1000)
    max_rev = st.sidebar.number_input("Max last 12 months revenues (USD)", min_value=0, value=10_000_000, step=1000)
    min_som = st.sidebar.number_input("Min market size or SOM (USD)", min_value=0, value=0, step=1000)
    max_som = st.sidebar.number_input("Max market size or SOM (USD)", min_value=0, value=100_000_000, step=1000)
    min_irr = st.sidebar.number_input("Min expected IRR (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.1)
    max_irr = st.sidebar.number_input("Max expected IRR (%)", min_value=0.0, max_value=100.0, value=100.0, step=0.1)

    folders = [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isdir(os.path.join(UPLOAD_FOLDER, f))]
    if not folders:
        st.info("No submissions yet.")
    else:

        # Build filtered records for dashboard metrics and charts
        records = []
        for fld in folders:
            base = os.path.join(UPLOAD_FOLDER, fld)
            # Read metadata
            info = {}
            meta_file = os.path.join(base, "info.txt")
            if os.path.exists(meta_file):
                for line in open(meta_file).read().splitlines():
                    if ":" in line:
                        k, v = line.split(":", 1)
                        info[k.strip()] = v.strip()
            # Read summary
            sum_path = os.path.join(base, "summary_gpt.txt")
            summary_dict = _parse_json_from_string(open(sum_path).read()) if os.path.exists(sum_path) else {}
            # Read status
            status = "Identified"
            spath = os.path.join(base, "status.json")
            if os.path.exists(spath):
                try:
                    status = json.load(open(spath)).get("status", status)
                except:
                    pass
            # Apply same filters as for display
            if hq_filter and hq_filter.lower() not in info.get("Country HQ", "").lower():
                continue
            if sector_filter and info.get("Sector", "") not in sector_filter:
                continue
            if main_country_filter and main_country_filter.lower() not in summary_dict.get("Main country of current operations", "").lower():
                continue
            sdg_raw = summary_dict.get("3 main SDGs targeted", "")
            sdg_list = [s.strip() for s in sdg_raw.split(";")] if isinstance(sdg_raw, str) else sdg_raw
            if sdg_filter and any(s not in sdg_list for s in sdg_filter):
                continue
            if geography_filter and summary_dict.get("Region of operation", "") not in geography_filter:
                continue
            if maturity_filter and summary_dict.get("Maturity stage", "") not in maturity_filter:
                continue
            # Numeric filters
            try:
                rev_val = float(summary_dict.get("Last 12 months revenues (USD)", "0").replace(",", ""))
            except:
                rev_val = 0.0
            if rev_val < min_rev or rev_val > max_rev:
                continue
            try:
                som_val = float(summary_dict.get("Market size or SOM (USD)", "0").replace(",", ""))
            except:
                som_val = 0.0
            if som_val < min_som or som_val > max_som:
                continue
            try:
                irr_val = float(summary_dict.get("Expected IRR (%)", "").replace("%", ""))
            except:
                irr_val = 0.0
            if irr_val < min_irr or irr_val > max_irr:
                continue
            # Last update time
            try:
                mtime = os.path.getmtime(os.path.join(base, "summary_gpt.txt"))
                last_update = datetime.fromtimestamp(mtime)
            except:
                last_update = None
            records.append({
                "Project": info.get("Project", fld),
                "Email": info.get("Email", ""),
                "Country HQ": info.get("Country HQ", ""),
                "Sector": info.get("Sector", ""),
                "Main country": summary_dict.get("Main country of current operations", ""),
                "SDGs": sdg_list,
                "Region of operation": summary_dict.get("Region of operation", ""),
                "Maturity": summary_dict.get("Maturity stage", ""),
                "Status": status,
                "Revenues": rev_val,
                "SOM": som_val,
                "IRR": irr_val,
                "LastUpdate": last_update,
            })
        # Display dashboard metrics and charts
        if records:
            df = pd.DataFrame(records)
            st.subheader("Dashboard Overview")
            # Key metrics at a glance
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("📁 Total Projects", len(df))
            col2.metric("💰 Avg Revenue (USD)", f"${df['Revenues'].mean():,.0f}")
            col3.metric("📈 Median IRR (%)", f"{df['IRR'].median():.1f}%")
            col4.metric("🌐 Unique Sectors", df["Sector"].nunique())
            st.markdown("---")  # separator before charts
            # Last project update
            latest = max(records, key=lambda x: x["LastUpdate"] or datetime.min)
            st.markdown(f"**Last project update:** {latest['LastUpdate'].strftime('%Y-%m-%d %H:%M:%S')} by {latest['Email']}")
            # Pie charts with titles
            col1, col2, col3 = st.columns(3)

            # Status distribution pie
            fig1, ax1 = plt.subplots(figsize=(4, 4))
            df["Status"].value_counts().plot.pie(
                autopct="%1.1f%%", ax=ax1, startangle=90
            )
            ax1.set_ylabel("")  # remove default ylabel
            ax1.set_title("By Status")
            col1.pyplot(fig1)

            # Region of operation distribution pie
            fig2, ax2 = plt.subplots(figsize=(4, 4))
            df["Region of operation"].value_counts().plot.pie(
                autopct="%1.1f%%", ax=ax2, startangle=90
            )
            ax2.set_ylabel("")
            ax2.set_title("By Region")
            col2.pyplot(fig2)

            # Sector distribution pie
            fig3, ax3 = plt.subplots(figsize=(4, 4))
            df["Sector"].value_counts().plot.pie(
                autopct="%1.1f%%", ax=ax3, startangle=90
            )
            ax3.set_ylabel("")
            ax3.set_title("By Sector")
            col3.pyplot(fig3)
            # Histogram per sector
            st.subheader("Submissions per Sector")
            st.bar_chart(df["Sector"].value_counts())

        for fld in sorted(folders):
            base = os.path.join(UPLOAD_FOLDER, fld)
            info = {}
            meta_file = os.path.join(base, "info.txt")
            if os.path.exists(meta_file):
                with open(meta_file) as mf:
                    for line in mf.read().splitlines():
                        if ":" in line:
                            k, v = line.split(":", 1)
                            info[k.strip()] = v.strip()

            # Apply Headquarters country filter
            if hq_filter:
                # hq_filter may be string or list
                if isinstance(hq_filter, list):
                    # treat as any-of: require at least one selected item in the field
                    if not any(item.lower() in info.get("Country HQ", "").lower() for item in hq_filter):
                        continue
                else:
                    if hq_filter.lower() not in info.get("Country HQ", "").lower():
                        continue
            if sector_filter and info.get("Sector", "") not in sector_filter:
                continue

            sum_path = os.path.join(base, "summary_gpt.txt")
            # Parse AI summary early for summary-based filters
            summary_dict = _parse_json_from_string(open(sum_path).read()) if os.path.exists(sum_path) else {}

            # Apply Main country of current operations filter
            if main_country_filter:
                if isinstance(main_country_filter, list):
                    if not any(item.lower() in summary_dict.get("Main country of current operations", "").lower() for item in main_country_filter):
                        continue
                else:
                    if main_country_filter.lower() not in summary_dict.get("Main country of current operations", "").lower():
                        continue

            # Apply SDG filter: fuzzy-matched summary list
            sdg_raw = summary_dict.get("3 main SDGs targeted", "")
            sdg_list = [s.strip() for s in sdg_raw.split(";")] if isinstance(sdg_raw, str) else sdg_raw
            if sdg_filter and any(s not in sdg_list for s in sdg_filter):
                continue

            # Apply geography filter
            if geography_filter and summary_dict.get("Region of operation", "") not in geography_filter:
                continue

            # Apply maturity filter
            if maturity_filter and summary_dict.get("Maturity stage", "") not in maturity_filter:
                continue

            # Apply numeric filters
            try:
                rev_val = float(summary_dict.get("Last 12 months revenues (USD)", "0").replace(",", ""))
            except:
                rev_val = 0.0
            if rev_val < min_rev or rev_val > max_rev:
                continue

            try:
                som_val = float(summary_dict.get("Market size or SOM (USD)", "0").replace(",", ""))
            except:
                som_val = 0.0
            if som_val < min_som or som_val > max_som:
                continue

            try:
                irr_val = float(summary_dict.get("Expected IRR (%)", "").replace("%", ""))
            except:
                irr_val = 0.0
            if irr_val < min_irr or irr_val > max_irr:
                continue

            status_path = os.path.join(base, "status.json")
            status_val = "Identified"
            if os.path.exists(status_path):
                try:
                    status_val = json.load(open(status_path)).get("status", status_val)
                except Exception:
                    pass

            with st.expander(f"{info.get('Project', fld)}"):
                st.markdown(f"**Project Name:** {info.get('Project', '–')}")
                st.markdown(f"**Contact Email:** {info.get('Email', '–')}")
                st.markdown(f"**Country:** {info.get('Country HQ', '–')}")
                st.markdown(f"**Sector:** {info.get('Sector', '–')}")

                st.markdown("**AI Summary:**")
                render_summary_grid(summary_dict)

                pdfs = [p for p in os.listdir(base) if p.lower().endswith(".pdf")]
                if pdfs:
                    pfile = pdfs[0]
                    with open(os.path.join(base, pfile), "rb") as pdf_bin:
                        st.download_button(
                            label=f"Download {pfile}",
                            data=pdf_bin,
                            file_name=pfile,
                            mime="application/pdf",
                            key=f"pdf_{fld}",        # ensure a unique key per project folder
                        )

                option = st.selectbox(
                    "Due Diligence / Operations stage",
                    REVIEW_STAGES,
                    index=REVIEW_STAGES.index(status_val) if status_val in REVIEW_STAGES else 0,
                    key=f"stage_{fld}",
                )
                if st.button("Save Status", key=f"save_{fld}"):
                    with open(status_path, "w") as sf:
                        json.dump({"status": option}, sf)
                    st.success("Status updated!")

    # ── EXPORT ───────────────────────────────────────────────
    rows = []
    for fld in folders:
        base = os.path.join(UPLOAD_FOLDER, fld)
        info = {}
        imeta = os.path.join(base, "info.txt")
        if os.path.exists(imeta):
            with open(imeta) as mf:
                for line in mf.read().splitlines():
                    if ":" in line:
                        k, v = line.split(":", 1)
                        info[k.strip()] = v.strip()
        status_val = json.load(open(os.path.join(base, "status.json")))['status'] if os.path.exists(os.path.join(base, "status.json")) else "Identified"
        parsed = _parse_json_from_string(open(os.path.join(base, "summary_gpt.txt")).read()) if os.path.exists(os.path.join(base, "summary_gpt.txt")) else {}
        row = {f: "" for f in AI_FIELDS}
        row.update({
            "Project": info.get("Project", fld),
            "Country HQ": info.get("Country HQ", ""),    # use HQ country
            "Sector": info.get("Sector", ""),
            "Status": status_val,
            "Email": info.get("Email", ""),
        })
        for k, v in parsed.items():
            if k in AI_FIELDS:
                row[k] = "; ".join(f"{m.get('Name','')}" if isinstance(m, dict) else str(m) for m in v) if isinstance(v, list) else v
        rows.append(row)
    if rows:
        df = pd.DataFrame(rows).reindex(columns=["Project", "Country", "Sector", "Status", "Email"] + AI_FIELDS)
        st.download_button("Download CSV", df.to_csv(index=False).encode(), "submissions.csv", "text/csv")
        buf = io.BytesIO()
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Subs")
        st.download_button("Download Excel", buf.getvalue(), "submissions.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


# ═════════════════════════════════════════════════════════════════
# ENTREPRENEUR DASHBOARD
# ═════════════════════════════════════════════════════════════════
else:
    st.title("📥 Impact Project Room")

    # Mode selector
    mode = st.sidebar.radio("Mode", ["Submit New", "Edit Existing"])

    if mode == "Edit Existing":
        st.subheader("Edit your existing submission")
        eid = st.text_input("Project ID")
        epin = st.text_input("PIN", type="password")
        if st.button("Load Submission"):
            cred_file = os.path.join(UPLOAD_FOLDER, eid or "", "credentials.json")
            if os.path.exists(cred_file):
                creds = json.load(open(cred_file))
                if creds.get("pin") == epin:
                    # load metadata & summary into session_state as before
                    base = os.path.join(UPLOAD_FOLDER, eid)
                    info = {}
                    for line in open(os.path.join(base, "info.txt")).read().splitlines():
                        if ":" in line:
                            k, v = line.split(":",1)
                            info[k.strip()] = v.strip()
                    summary = {}
                    try:
                        summary = json.loads(open(os.path.join(base, "summary_gpt.txt")).read())
                    except:
                        pass
                    st.session_state.form_meta = {
                        "Project": info.get("Project",""),
                        "Incorporation date": info.get("Incorporation date", datetime.now().date()),
                        "Country HQ": info.get("Country HQ",""),
                        "Sector": info.get("Sector",""),
                        "Email": info.get("Email",""),
                    }
                    st.session_state.form_files = []
                    st.session_state.form_summ = summary
                    st.session_state.edit_folder = eid
                    st.session_state.stage = "review"
                    _rerun()
                else:
                    st.error("Invalid PIN.")
            else:
                st.error("Invalid Project ID.")

    # ---------- two-stage form -------------
    if "stage" not in st.session_state:
        st.session_state.stage = "input"

    # ── Stage 1: basic inputs + AI generation
    if mode == "Submit New" and st.session_state.stage == "input":
        # Mutual NDA download button (outside the form)
        nda_path = Path(__file__).with_name("Mutual agreement.pdf")
        if nda_path.exists():
            with nda_path.open("rb") as nda_file:
                st.download_button(
                    label="Download Mutual NDA (PDF)",
                    data=nda_file.read(),
                    file_name="Mutual_agreement.pdf",
                    mime="application/pdf",
                    key="nda_download",
                )
        else:
            st.warning("NDA PDF not found.")

        with st.form("project_form"):
            proj     = st.text_input("Project registered name")
            inc_date = st.date_input("Date of incorporation")
            sector   = st.selectbox("Primary sector / theme", SECTOR_OPTIONS)
            country  = st.selectbox("Headquarters country", COUNTRY_OPTIONS)
            files    = st.file_uploader(
                "Upload up to 5 files (any format)",
                accept_multiple_files=True,
            )
            email    = st.text_input("Contact e-mail")

            nda = st.checkbox("I accept the NDA")

            generate_btn = st.form_submit_button("Generate AI Summary")

        if generate_btn:
            # validation
            if not (proj and email and nda and files):
                st.warning("Fill required fields, upload at least one file, accept NDA.")
                st.stop()
            if len(files) > 5:
                st.warning("Max 5 files.")
                st.stop()

            # run AI on first uploaded file (PDF, DOCX, PPTX, XLSX)
            from tempfile import NamedTemporaryFile

            first_file = files[0] if files else None
            summary = {}
            if first_file:
                suffix = Path(first_file.name).suffix
                with NamedTemporaryFile(delete=True, suffix=suffix) as tmp:
                    tmp.write(first_file.read())
                    tmp.flush()
                    text = extract_text_from_file(tmp.name)
                summary = summarize_project_with_gpt(text)

            # stash in session state
            st.session_state.form_meta  = {
                "Project": proj,
                "Incorporation date": inc_date,
                "Country HQ": country,
                "Sector": sector,
                "Email": email,
            }
            st.session_state.form_files = files
            st.session_state.form_summ  = summary or {k: "Unknown" for k in AI_FIELDS}
            st.session_state.stage = "review"
            _rerun()

    # ── Stage 2: review AI output, edit, confirm
    if st.session_state.stage == "review":
        # ── Review & edit AI-generated summary with typed inputs
        st.header("Review & edit AI-generated summary")
        edited = {}

        # Project Name and Business Model: free text
        edited["Project Name"] = st.text_input(
            "Project Name",
            value=st.session_state.form_summ.get("Project Name", ""),
            key="edit_Project Name",
        )
        edited["Business Model"] = st.text_input(
            "Business Model",
            value=st.session_state.form_summ.get("Business Model", ""),
            key="edit_Business Model",
        )

        # Maturity stage: select from standard list with fallback for custom values
        # Graceful fallback if AI returned a custom stage not in list
        ms_val = st.session_state.form_summ.get("Maturity stage", "")
        default_index = MATURITY_STAGES.index(ms_val) if ms_val in MATURITY_STAGES else 0
        edited["Maturity stage"] = st.selectbox(
            "Maturity stage",
            MATURITY_STAGES,
            index=default_index,
            key="edit_Maturity stage",
        )

        # Core Team, Impact Area, Key risks, Barrier(s) to entry
        for field in ["Core team", "Impact Area", "Key risks", "Barrier(s) to entry"]:
            edited[field] = st.text_input(
                field,
                value=st.session_state.form_summ.get(field, ""),
                key=f"edit_{field}",
            )

        # Numeric fields: integer / float inputs
        # Last 12 months revenues
        rev_val_raw = st.session_state.form_summ.get("Last 12 months revenues (USD)", "")
        rev_raw = str(rev_val_raw).replace(",", "")
        rev_default = int(rev_raw) if rev_raw.isdigit() else 0
        edited["Last 12 months revenues (USD)"] = st.number_input(
            "Last 12 months revenues (USD)",
            min_value=0,
            value=rev_default,
            step=1,
            format="%d",
            key="edit_Last 12 months revenues",
        )

        # Market size or SOM
        som_val_raw = st.session_state.form_summ.get("Market size or SOM (USD)", "")
        som_raw = str(som_val_raw).replace(",", "")
        som_default = int(som_raw) if som_raw.isdigit() else 0
        edited["Market size or SOM (USD)"] = st.number_input(
            "Market size or SOM (USD)",
            min_value=0,
            value=som_default,
            step=1,
            format="%d",
            key="edit_Market size or SOM",
        )

        # Financing need or round size
        fn_val_raw = st.session_state.form_summ.get("Financing need or round size (USD)", "")
        fn_raw = str(fn_val_raw).replace(",", "")
        fn_default = int(fn_raw) if fn_raw.isdigit() else 0
        edited["Financing need or round size (USD)"] = st.number_input(
            "Financing need or round size (USD)",
            min_value=0,
            value=fn_default,
            step=1,
            format="%d",
            key="edit_Financing need or round size",
        )

        # Breakeven year
        by_raw = st.session_state.form_summ.get("Breakeven year", "")
        by_default = int(by_raw) if str(by_raw).isdigit() else datetime.now().year
        edited["Breakeven year"] = st.number_input(
            "Breakeven year",
            min_value=1900,
            max_value=2100,
            value=by_default,
            step=1,
            key="edit_Breakeven year",
        )

        # Expected IRR as percentage
        irr_val_raw = st.session_state.form_summ.get("Expected IRR (%)", "")
        irr_raw = str(irr_val_raw).replace("%", "")
        try:
            irr_default = float(irr_raw)
        except:
            irr_default = 0.0
        edited["Expected IRR (%)"] = st.number_input(
            "Expected IRR (%)",
            min_value=0.0,
            max_value=100.0,
            value=irr_default,
            step=0.1,
            format="%.2f",
            key="edit_Expected IRR",
        )

        # Use of proceeds
        uop_val_raw = st.session_state.form_summ.get("Use of proceeds (%)", "")
        uop_raw = str(uop_val_raw).replace("%", "")
        try:
            uop_default = float(uop_raw)
        except:
            uop_default = 0.0
        edited["Use of proceeds (%)"] = st.number_input(
            "Use of proceeds (%)",
            min_value=0.0,
            max_value=100.0,
            value=uop_default,
            step=1.0,
            format="%.1f",
            key="edit_Use of proceeds",
        )

        # SDGs targeted: up to 3 selections, with fallback for non-standard values
        raw_sdgs = st.session_state.form_summ.get("3 main SDGs targeted", "")
        if isinstance(raw_sdgs, str):
            raw_list = [s.strip() for s in raw_sdgs.split(";") if s.strip()]
        else:
            raw_list = raw_sdgs

        # Fuzzy-match AI output to canonical SDG options
        default_sdgs = _match_sdgs(raw_list)

        edited["3 main SDGs targeted"] = st.multiselect(
            "3 main SDGs targeted",
            SDG_OPTIONS,
            default=default_sdgs,
            max_selections=3,
            key="edit_3 main SDGs targeted",
        )

        # Other free-text fields
        for field in ["Problem", "Solution"]:
            edited[field] = st.text_area(
                field,
                value=st.session_state.form_summ.get(field, ""),
                key=f"edit_{field}",
            )

        if st.button("✅ Confirm & Submit"):
            if "edit_folder" in st.session_state:
                _update_submission(
                    st.session_state.edit_folder,
                    st.session_state.form_meta,
                    edited,
                )
                fld = os.path.join(UPLOAD_FOLDER, st.session_state.edit_folder)
                token = None
            else:
                fld, token = _save_submission(
                    st.session_state.form_meta,
                    st.session_state.form_files,
                    edited,
                )
                st.session_state._last_submission_token = token
            # Send email alert
            subject = "📥 New Impact Project Room submission"
            body = (
                f"Project: {edited.get('Project Name', st.session_state.form_meta['Project'])}\n"
                f"Sector:  {st.session_state.form_meta['Sector']}\n"
                f"Country: {st.session_state.form_meta['Country HQ']}\n"
                f"Uploaded: {datetime.now().isoformat(timespec='seconds')}"
            )
            email_admin(subject, body)
            st.session_state.submitted_details = edited
            # Save for next stage
            st.session_state._last_submission_fld = fld
            # Move to completion stage
            st.session_state.stage = "done"
            _rerun()

    # ── Stage 3: submission complete page  ------------------------
    if st.session_state.stage == "done":
        st.title("🎉 Submission Complete")
        st.write("Thank you! Your project has been submitted and is now awaiting review by our team.")
        st.write("We will notify you at ", st.session_state.form_meta.get("Email"), " once we have an update.")
        # Optionally display a summary of what was submitted:
        st.subheader("Your submitted details:")
        for field, val in st.session_state.submitted_details.items():
            st.markdown(f"- **{field}:** {val}")
        # Display Project ID and PIN on new submissions
        if not st.session_state.get("edit_folder"):
            st.subheader("Your Project ID and PIN:")
            proj_id = os.path.basename(st.session_state._last_submission_fld)
            pin = st.session_state._last_submission_token
            st.markdown(f"- **Project ID:** `{proj_id}`")
            st.markdown(f"- **PIN:** `{pin}`")
            st.info("Please save your Project ID and PIN to edit your submission later.")
        # Offer a “Start over” button to clear state
        if st.button("Submit another project"):
            for k in ("stage", "form_meta", "form_files", "form_summ", "submitted_details", "_last_submission_fld", "_last_submission_token", "edit_folder"):
                st.session_state.pop(k, None)
            _rerun()
